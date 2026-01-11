"""
PPTX formatter for presentations.

Generates PowerPoint presentations with:
- Title slide
- Outline/agenda slide
- Content slides (one concept per slide)
- Conclusion slide
- References slide
- Theme-based color schemes
"""

from pathlib import Path
from typing import Any, Optional
from datetime import datetime

from shared.formatters.base_formatter import BaseFormatter
from shared.utils.metadata_manager import create_presentation_metadata

# Conditional import for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Theme color schemes
THEMES = {
    "stem": {
        "name": "Modern Tech Blue",
        "primary": "2E5C8A",
        "secondary": "4A90E2",
        "accent": "FF6B35",
        "text": "333333",
        "background": "FFFFFF",
    },
    "sciences": {
        "name": "Scientific Green",
        "primary": "2D5A27",
        "secondary": "4A8C41",
        "accent": "FFD700",
        "text": "333333",
        "background": "FFFFFF",
    },
    "humanities": {
        "name": "Classic Warm",
        "primary": "8B4513",
        "secondary": "CD853F",
        "accent": "DAA520",
        "text": "333333",
        "background": "FFF8F0",
    },
    "business": {
        "name": "Corporate Blue",
        "primary": "1C3D5A",
        "secondary": "4682B4",
        "accent": "FF8C00",
        "text": "333333",
        "background": "FFFFFF",
    },
    "default": {
        "name": "Modern Minimalist",
        "primary": "2E5C8A",
        "secondary": "6B7B8C",
        "accent": "E74C3C",
        "text": "333333",
        "background": "FFFFFF",
    },
}


class PptxFormatter(BaseFormatter):
    """Formatter for generating PowerPoint presentations."""

    @property
    def material_type(self) -> str:
        return "presentations"

    @property
    def output_format(self) -> str:
        return "pptx"

    def __init__(self, subject: str, topic: str, theme: str = "default"):
        super().__init__(subject, topic)
        self.theme_name = theme
        self.theme = THEMES.get(theme, THEMES["default"])

    def generate(self, content: dict[str, Any]) -> Path:
        """
        Generate PowerPoint presentation.

        Args:
            content: Dictionary containing:
                - version: str
                - theme: str (optional)
                - outline: list[str] (main topics)
                - slides: list[dict] with:
                    - title: str
                    - bullets: list[str]
                    - notes: str (optional)
                - conclusion: dict with key_points
                - references: list[dict]
                - update_highlights: str (optional)

        Returns:
            Path to the generated PowerPoint file
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PPTX generation. "
                "Install with: pip install python-pptx"
            )

        # Update theme if specified
        if content.get("theme"):
            self.theme_name = content["theme"]
            self.theme = THEMES.get(content["theme"], THEMES["default"])

        self.ensure_output_directory()
        output_path = self.get_output_path()

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs, content)

        # Update highlights (if updating)
        if content.get("update_highlights"):
            self._add_update_slide(prs, content)

        # Outline slide
        if content.get("outline"):
            self._add_outline_slide(prs, content)

        # Content slides
        for slide_data in content.get("slides", []):
            self._add_content_slide(prs, slide_data)

        # Conclusion slide
        if content.get("conclusion"):
            self._add_conclusion_slide(prs, content)

        # References slide
        if content.get("references"):
            self._add_references_slide(prs, content)

        # Save presentation
        prs.save(str(output_path))

        # Save metadata
        metadata = create_presentation_metadata(
            topic=self.topic,
            subject=self.subject,
            number_of_slides=len(prs.slides),
            theme={
                "type": "auto_selected" if not content.get("theme") else "user_selected",
                "name": self.theme["name"],
                "primary_color": f"#{self.theme['primary']}",
                "secondary_color": f"#{self.theme['secondary']}",
                "accent_color": f"#{self.theme['accent']}",
            },
            reference=content.get("reference"),
            features=content.get("features", ["Academic tone", "Modern theme"])
        )
        metadata["current_version"] = content.get("version", "v1.0")
        self.save_metadata(metadata)

        return output_path

    def _hex_to_rgb(self, hex_color: str) -> RgbColor:
        """Convert hex color to RgbColor."""
        hex_color = hex_color.lstrip('#')
        return RgbColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def _add_title_slide(self, prs: Presentation, content: dict[str, Any]):
        """Add title slide."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Background color (optional)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(self.theme["background"])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = self.topic
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(self.theme["primary"])
        title_para.alignment = PP_ALIGN.CENTER

        # Subject
        subject_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(11.333), Inches(0.5)
        )
        subject_frame = subject_box.text_frame
        subject_para = subject_frame.paragraphs[0]
        subject_para.text = self.subject
        subject_para.font.size = Pt(24)
        subject_para.font.color.rgb = self._hex_to_rgb(self.theme["secondary"])
        subject_para.alignment = PP_ALIGN.CENTER

        # Date and version
        version = content.get("version", "v1.0")
        meta_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(11.333), Inches(0.5)
        )
        meta_frame = meta_box.text_frame
        meta_para = meta_frame.paragraphs[0]
        meta_para.text = f"{self._format_date()}  |  {version}"
        meta_para.font.size = Pt(14)
        meta_para.font.color.rgb = self._hex_to_rgb(self.theme["text"])
        meta_para.alignment = PP_ALIGN.CENTER

    def _add_outline_slide(self, prs: Presentation, content: dict[str, Any]):
        """Add outline/agenda slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        self._add_slide_title(slide, "Outline")

        # Outline items
        outline = content.get("outline", [])
        y_position = 1.8

        for i, item in enumerate(outline, 1):
            text_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_position), Inches(10), Inches(0.5)
            )
            text_frame = text_box.text_frame
            para = text_frame.paragraphs[0]
            para.text = f"{i}. {item}"
            para.font.size = Pt(20)
            para.font.color.rgb = self._hex_to_rgb(self.theme["text"])
            y_position += 0.6

    def _add_content_slide(self, prs: Presentation, slide_data: dict[str, Any]):
        """Add a content slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title = slide_data.get("title", "Content")
        self._add_slide_title(slide, title)

        # Bullets
        bullets = slide_data.get("bullets", [])
        y_position = 1.8

        for bullet in bullets[:7]:  # Max 7 bullets per slide
            text_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_position), Inches(10), Inches(0.5)
            )
            text_frame = text_box.text_frame
            para = text_frame.paragraphs[0]
            para.text = f"• {bullet}"
            para.font.size = Pt(18)
            para.font.color.rgb = self._hex_to_rgb(self.theme["text"])
            y_position += 0.6

        # Add notes if present
        notes = slide_data.get("notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_update_slide(self, prs: Presentation, content: dict[str, Any]):
        """Add update highlights slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        version = content.get("version", "v1.1")
        self._add_slide_title(slide, f"UPDATE HIGHLIGHTS - {version}")

        highlights = content.get("update_highlights", "")

        text_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(1.8), Inches(10), Inches(4)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        para = text_frame.paragraphs[0]
        para.text = highlights
        para.font.size = Pt(16)
        para.font.color.rgb = self._hex_to_rgb(self.theme["text"])

    def _add_conclusion_slide(self, prs: Presentation, content: dict[str, Any]):
        """Add conclusion slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Key Takeaways")

        conclusion = content.get("conclusion", {})
        key_points = conclusion.get("key_points", [])
        y_position = 1.8

        for point in key_points[:5]:
            text_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_position), Inches(10), Inches(0.5)
            )
            text_frame = text_box.text_frame
            para = text_frame.paragraphs[0]
            para.text = f"✓ {point}"
            para.font.size = Pt(20)
            para.font.color.rgb = self._hex_to_rgb(self.theme["accent"])
            y_position += 0.7

    def _add_references_slide(self, prs: Presentation, content: dict[str, Any]):
        """Add references slide."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "References")

        refs = content.get("references", [])
        y_position = 1.8

        for i, ref in enumerate(refs, 1):
            ref_text = self._format_reference(ref)
            text_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_position), Inches(10), Inches(0.5)
            )
            text_frame = text_box.text_frame
            para = text_frame.paragraphs[0]
            para.text = f"{i}. {ref_text}"
            para.font.size = Pt(14)
            para.font.color.rgb = self._hex_to_rgb(self.theme["text"])
            y_position += 0.5

    def _add_slide_title(self, slide, title: str):
        """Add title to a slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12.333), Inches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(self.theme["primary"])

    def _format_reference(self, ref: dict[str, Any]) -> str:
        """Format a reference for display."""
        ref_type = ref.get("type", "general")
        ref_content = ref.get("content", "General knowledge")

        if ref_type == "book":
            return f"{ref_content} (Book)"
        elif ref_type == "url":
            return ref_content
        elif ref_type == "notes_skill":
            return f"Notes: {ref_content}"
        else:
            return ref_content
